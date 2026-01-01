from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, Any, List
import io
import json


class PPTXGenerator:
    """Generate professional PPTX presentations following Indian engineering college standards"""
    
    def _ensure_dict(self, data: Any, default: Any = None) -> Dict:
        """Ensure data is a dictionary, parsing JSON if needed"""
        if data is None:
            return default if default is not None else {}
        if isinstance(data, str):
            try:
                parsed = json.loads(data)
                return parsed if isinstance(parsed, dict) else (default if default is not None else {})
            except json.JSONDecodeError:
                return default if default is not None else {}
        if isinstance(data, dict):
            return data
        return default if default is not None else {}
    
    def _ensure_list(self, data: Any) -> list:
        """Ensure data is a list"""
        if data is None:
            return []
        if isinstance(data, list):
            return data
        if isinstance(data, str):
            return [data]
        return []
    
    def _add_bullet_slide(self, prs: Presentation, title: str, bullets: List[str]):
        """Add a slide with title and bullet points"""
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title_shape = slide.shapes.title
        body = slide.placeholders[1]
        
        title_shape.text = title
        tf = body.text_frame
        tf.clear()
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = str(bullet)
            p.level = 0
        
        return slide
    
    def generate_slides(self, project_data: Dict[str, Any]) -> bytes:
        """
        Generate professional PPTX presentation following Indian engineering standards
        
        Args:
            project_data: Project JSON from LLM
            
        Returns:
            PPTX file as bytes
        """
        project_data = self._ensure_dict(project_data)
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)
        
        # ============ SLIDE 1: TITLE SLIDE ============
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = project_data.get('title', 'Project Presentation')
        subtitle.text = f"""A Semester Project Presentation
        
Difficulty: {project_data.get('difficulty', 'Intermediate')} | Timeline: {project_data.get('timeline_days', 45)} days

[University Name / College Name]
[Student Name(s)]
[Guide Name]"""
        
        # ============ SLIDE 2: ABSTRACT ============
        abstract = project_data.get('abstract', 'Project abstract to be added.')
        if len(abstract) > 400:
            abstract = abstract[:400] + "..."
        self._add_bullet_slide(prs, "Abstract", [abstract])
        
        # ============ SLIDE 3: PROBLEM STATEMENT ============
        introduction = self._ensure_dict(project_data.get('introduction', {}))
        if isinstance(project_data.get('introduction'), str):
            self._add_bullet_slide(prs, "Problem Statement", [project_data.get('introduction', '')])
        else:
            problem_bullets = []
            if introduction.get('problem_statement'):
                problem_bullets.append(introduction['problem_statement'])
            if introduction.get('motivation'):
                problem_bullets.append(f"Motivation: {introduction['motivation']}")
            if problem_bullets:
                self._add_bullet_slide(prs, "Problem Statement", problem_bullets)
        
        # ============ SLIDE 4: OBJECTIVES ============
        objectives = self._ensure_list(project_data.get('objectives', []))
        if objectives:
            self._add_bullet_slide(prs, "Objectives", objectives[:6])
        
        # ============ SLIDE 5: SCOPE ============
        scope = project_data.get('scope', {})
        scope_bullets = []
        if isinstance(scope, str):
            scope_bullets = [scope]
        elif isinstance(scope, dict):
            if scope.get('in_scope'):
                for item in self._ensure_list(scope.get('in_scope', []))[:3]:
                    scope_bullets.append(f"✓ {item}")
            if scope.get('target_users'):
                scope_bullets.append(f"Target Users: {scope.get('target_users')}")
        if scope_bullets:
            self._add_bullet_slide(prs, "Scope", scope_bullets)
        
        # ============ SLIDE 6: LITERATURE SURVEY ============
        literature = self._ensure_list(project_data.get('literature_survey', []))
        if literature:
            lit_bullets = []
            for paper in literature[:4]:
                if isinstance(paper, dict):
                    lit_bullets.append(f"{paper.get('paper_title', 'Paper')} ({paper.get('year', 'N/A')})")
                elif isinstance(paper, str):
                    lit_bullets.append(paper[:80] + "..." if len(paper) > 80 else paper)
            if lit_bullets:
                self._add_bullet_slide(prs, "Literature Survey", lit_bullets)
        
        # ============ SLIDE 7: EXISTING VS PROPOSED SYSTEM ============
        system_study = self._ensure_dict(project_data.get('system_study', {}))
        existing = self._ensure_dict(system_study.get('existing_system', {}))
        proposed = self._ensure_dict(system_study.get('proposed_system', {}))
        
        if existing or proposed:
            comparison_bullets = []
            if existing.get('limitations'):
                comparison_bullets.append("Existing System Limitations:")
                for lim in self._ensure_list(existing.get('limitations', []))[:2]:
                    comparison_bullets.append(f"  ✗ {lim}")
            if proposed.get('advantages'):
                comparison_bullets.append("Proposed System Advantages:")
                for adv in self._ensure_list(proposed.get('advantages', []))[:2]:
                    comparison_bullets.append(f"  ✓ {adv}")
            if comparison_bullets:
                self._add_bullet_slide(prs, "Existing vs Proposed System", comparison_bullets)
        
        # ============ SLIDE 8: SYSTEM ARCHITECTURE ============
        system_design = self._ensure_dict(project_data.get('system_design', {}))
        arch = self._ensure_dict(system_design.get('architecture', {}))
        if arch:
            arch_bullets = []
            if arch.get('type'):
                arch_bullets.append(f"Architecture: {arch.get('type')}")
            if arch.get('components'):
                arch_bullets.append("Components:")
                for comp in self._ensure_list(arch.get('components', []))[:4]:
                    arch_bullets.append(f"  • {comp}")
            if arch.get('diagram_description'):
                arch_bullets.append(f"[Diagram: {arch.get('diagram_description')[:100]}...]")
            if arch_bullets:
                self._add_bullet_slide(prs, "System Architecture", arch_bullets)
        
        # ============ SLIDE 9: TECHNOLOGY STACK ============
        tech_stack = self._ensure_dict(project_data.get('technology_stack', {}))
        if tech_stack:
            tech_bullets = []
            
            if isinstance(tech_stack.get('frontend'), dict):
                frontend = tech_stack.get('frontend', {})
                tech_bullets.append(f"Frontend: {', '.join(self._ensure_list(frontend.get('technologies', [])))}")
            elif tech_stack.get('frontend'):
                tech_bullets.append(f"Frontend: {tech_stack.get('frontend')}")
            
            if isinstance(tech_stack.get('backend'), dict):
                backend = tech_stack.get('backend', {})
                tech_bullets.append(f"Backend: {', '.join(self._ensure_list(backend.get('technologies', [])))}")
            elif tech_stack.get('backend'):
                tech_bullets.append(f"Backend: {tech_stack.get('backend')}")
            
            if isinstance(tech_stack.get('database'), dict):
                database = tech_stack.get('database', {})
                tech_bullets.append(f"Database: {database.get('type', 'N/A')}")
            elif tech_stack.get('database'):
                tech_bullets.append(f"Database: {tech_stack.get('database')}")
            
            if tech_stack.get('other_tools'):
                tech_bullets.append(f"Tools: {', '.join(self._ensure_list(tech_stack.get('other_tools', [])))}")
            
            if tech_bullets:
                self._add_bullet_slide(prs, "Technology Stack", tech_bullets)
        
        # ============ SLIDE 10: PROJECT MODULES ============
        modules = self._ensure_list(project_data.get('modules', []))
        if modules:
            module_bullets = []
            for module in modules[:6]:
                if isinstance(module, str):
                    module_bullets.append(module)
                elif isinstance(module, dict):
                    weeks = module.get('duration_weeks', '')
                    week_str = f" ({weeks} weeks)" if weeks else ""
                    module_bullets.append(f"{module.get('name', 'Module')}{week_str}")
            
            self._add_bullet_slide(prs, "Project Modules", module_bullets)
        
        # ============ SLIDE 11: DATABASE DESIGN ============
        db_design = self._ensure_dict(project_data.get('database_design', {}))
        if db_design:
            db_bullets = []
            if db_design.get('database_type'):
                db_bullets.append(f"Database Type: {db_design.get('database_type')}")
            if db_design.get('normalization'):
                db_bullets.append(f"Normalization: {db_design.get('normalization')[:60]}...")
            
            tables = self._ensure_list(db_design.get('tables', []))
            if tables:
                db_bullets.append("Tables:")
                for table in tables[:4]:
                    if isinstance(table, dict):
                        db_bullets.append(f"  • {table.get('name', 'table')}")
            
            if db_bullets:
                self._add_bullet_slide(prs, "Database Design", db_bullets)
        
        # ============ SLIDE 12: IMPLEMENTATION HIGHLIGHTS ============
        impl = self._ensure_dict(project_data.get('implementation', {}))
        if impl and isinstance(impl, dict):
            impl_bullets = []
            if impl.get('development_environment'):
                impl_bullets.append(f"Environment: {impl.get('development_environment')}")
            if impl.get('coding_standards'):
                impl_bullets.append(f"Standards: {impl.get('coding_standards')}")
            
            key_algos = self._ensure_list(impl.get('key_algorithms', []))
            if key_algos:
                impl_bullets.append("Key Algorithms:")
                for algo in key_algos[:3]:
                    if isinstance(algo, dict):
                        impl_bullets.append(f"  • {algo.get('name', 'Algorithm')}")
                    elif isinstance(algo, str):
                        impl_bullets.append(f"  • {algo}")
            
            if impl_bullets:
                self._add_bullet_slide(prs, "Implementation", impl_bullets)
        
        # ============ SLIDE 13: TESTING ============
        testing = self._ensure_dict(project_data.get('testing', {}))
        if testing or isinstance(project_data.get('testing'), str):
            test_bullets = []
            
            if isinstance(project_data.get('testing'), str):
                test_bullets.append(project_data.get('testing'))
            else:
                if testing.get('testing_methodology'):
                    test_bullets.append(f"Methodology: {testing.get('testing_methodology')}")
                if testing.get('tools_used'):
                    test_bullets.append(f"Tools: {', '.join(self._ensure_list(testing.get('tools_used', [])))}")
                
                test_summary = self._ensure_dict(testing.get('test_summary', {}))
                if test_summary:
                    test_bullets.append(f"Total Test Cases: {test_summary.get('total_cases', 'N/A')}")
                    test_bullets.append(f"Pass Rate: {test_summary.get('pass_percentage', 'N/A')}")
            
            if test_bullets:
                self._add_bullet_slide(prs, "Testing", test_bullets)
        
        # ============ SLIDE 14: SCREENSHOTS / DEMO ============
        screenshots = self._ensure_list(project_data.get('screenshots', []))
        if screenshots:
            screen_bullets = []
            for screen in screenshots[:5]:
                if isinstance(screen, dict):
                    screen_bullets.append(f"📱 {screen.get('screen_name', 'Screen')}: {screen.get('description', '')[:50]}...")
                elif isinstance(screen, str):
                    screen_bullets.append(f"📱 {screen[:60]}...")
            
            if screen_bullets:
                self._add_bullet_slide(prs, "Screenshots / Demo", screen_bullets)
        
        # ============ SLIDE 15: CONCLUSION ============
        conclusion = project_data.get('conclusion', '')
        if conclusion:
            if len(conclusion) > 400:
                conclusion = conclusion[:400] + "..."
            self._add_bullet_slide(prs, "Conclusion", [conclusion])
        
        # ============ SLIDE 16: FUTURE SCOPE ============
        future_scope = project_data.get('future_scope', [])
        if future_scope:
            if isinstance(future_scope, str):
                future_bullets = [future_scope]
            else:
                future_bullets = self._ensure_list(future_scope)[:5]
            
            self._add_bullet_slide(prs, "Future Scope", future_bullets)
        
        # ============ SLIDE 17: VIVA PREPARATION ============
        viva_questions = self._ensure_list(project_data.get('viva_questions', []))
        if viva_questions:
            viva_bullets = []
            for q in viva_questions[:5]:
                if isinstance(q, str):
                    viva_bullets.append(f"Q: {q[:80]}...")
                elif isinstance(q, dict):
                    question = q.get('question', q.get('q', ''))
                    if question:
                        viva_bullets.append(f"Q: {question[:80]}...")
            
            if viva_bullets:
                self._add_bullet_slide(prs, "Sample Viva Questions", viva_bullets)
        
        # ============ SLIDE 18: REFERENCES ============
        references = self._ensure_list(project_data.get('references', []))
        if references:
            ref_bullets = []
            for idx, ref in enumerate(references[:5], 1):
                if isinstance(ref, str):
                    ref_bullets.append(f"[{idx}] {ref[:70]}...")
                elif isinstance(ref, dict):
                    citation = ref.get('citation', '')
                    ref_bullets.append(f"[{idx}] {citation[:70]}...")
            
            if ref_bullets:
                self._add_bullet_slide(prs, "References", ref_bullets)
        
        # ============ SLIDE 19: THANK YOU ============
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Thank You!"
        subtitle.text = """Questions?

[Student Name(s)]
[Email / Contact]

Guide: [Guide Name]
[Department Name]
[College/University Name]"""
        
        # Save to bytes
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        return file_stream.read()


# Singleton instance
pptx_generator = PPTXGenerator()
